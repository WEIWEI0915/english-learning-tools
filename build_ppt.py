#!/usr/bin/env python3
"""
English Singsing 英语教学PPT生成器
基于动画视频的互动教学模式：看视频 → 学词汇 → 练对话 → 角色扮演 → 词卡扩展
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn, nsmap
import os
import copy
from pathlib import Path

# ============================================================
# CONFIGURATION
# ============================================================

VIDEO_DIR = Path(__file__).parent / "videos"
OUTPUT_FILE = Path(__file__).parent / "英语启蒙教学课件-English_Singsing版.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_EN = 'Comic Sans MS'
FONT_CN = '微软雅黑'
FONT_TITLE = '微软雅黑'

# Unit color schemes (primary, secondary, accent)
UNIT_THEMES = {
    1: ('156082', '0F9ED5', 'E97132'),  # Ocean Blue + Orange
    2: ('E97132', 'FFB347', '156082'),  # Orange + Blue
    3: ('196B24', '4EA72E', '0F9ED5'),  # Green + Sky
    4: ('A02B93', 'C77DFF', 'E97132'),  # Purple + Orange
    5: ('0F9ED5', '7EC8E3', '156082'),  # Sky Blue + Navy
    6: ('E97132', 'FF9F43', '196B24'),  # Orange + Green
    7: ('4EA72E', '7ED957', 'E97132'),  # Lime + Orange
    8: ('156082', '3A8FC4', 'A02B93'),  # Navy + Purple
    9: ('0F9ED5', '5DBBDE', '196B24'),  # Sky + Green
    10: ('A02B93', 'D08BFF', '0F9ED5'), # Purple + Sky
}

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)

# ============================================================
# CONTENT DATA
# ============================================================

UNITS = [
    {
        'num': 1,
        'title': 'Greetings 问候',
        'subtitle': "Hello! How are you?",
        'icon': '👋',
        'vocab': [
            ('hello', '你好'), ('goodbye', '再见'), ('morning', '早晨'),
            ('afternoon', '下午'), ('evening', '傍晚'), ('night', '晚上'),
            ('friend', '朋友'), ('name', '名字'), ('nice', '好的'),
            ('meet', '遇见'), ('thank you', '谢谢'), ('welcome', '欢迎'),
        ],
        'dialogue': [
            ('A', 'Good morning!', '早上好！'),
            ('B', 'Good morning! How are you?', '早上好！你好吗？'),
            ('A', "I'm fine, thank you. And you?", '我很好，谢谢。你呢？'),
            ('B', "I'm great! Nice to meet you.", '我很好！很高兴认识你。'),
            ('A', 'Nice to meet you, too.', '我也很高兴认识你。'),
            ('B', "What's your name?", '你叫什么名字？'),
            ('A', "My name is Tom. What's your name?", '我叫Tom。你叫什么？'),
            ('B', "My name is Lily.", '我叫Lily。'),
        ],
        'grammar': [
            ('Be动词 (am/is/are)', 'I am / You are / He is / She is'),
            ('问候句型', 'Good morning/afternoon/evening!'),
            ('问候句型', 'How are you? → I\'m fine/great/good.'),
        ],
        'extended_vocab': [
            ('hi', '嗨'), ('hey', '嘿'), ('please', '请'),
            ('sorry', '对不起'), ('excuse me', '打扰一下'),
            ('see you', '再见'), ('take care', '保重'), ('good night', '晚安'),
        ],
        'video_search': 'English Singsing Greeting Good morning How are you',
        'video_bv': 'BV1L4411h77Q',  # 英语对话[问候] Good morning. How are you
        'video_file': 'greetings_dialogue.mp4',
    },
    {
        'num': 2,
        'title': 'Colors & Shapes 颜色与形状',
        'subtitle': "What color is it?",
        'icon': '🎨',
        'vocab': [
            ('red', '红色'), ('blue', '蓝色'), ('yellow', '黄色'),
            ('green', '绿色'), ('orange', '橙色'), ('purple', '紫色'),
            ('pink', '粉色'), ('white', '白色'), ('black', '黑色'),
            ('circle', '圆形'), ('square', '正方形'), ('triangle', '三角形'),
        ],
        'dialogue': [
            ('A', 'What color is it?', '这是什么颜色？'),
            ('B', "It's red!", '这是红色！'),
            ('A', 'What color do you like?', '你喜欢什么颜色？'),
            ('B', 'I like blue. What about you?', '我喜欢蓝色。你呢？'),
            ('A', 'I like yellow. It\'s sunny!', '我喜欢黄色。像阳光一样！'),
            ('B', "Look! A rainbow! Red, orange, yellow...", '看！彩虹！红色、橙色、黄色……'),
            ('A', 'Green, blue, purple! So beautiful!', '绿色、蓝色、紫色！太美了！'),
        ],
        'grammar': [
            ('What 疑问句', 'What color is it? → It\'s + 颜色.'),
            ('喜好表达', 'I like + 颜色/名词.'),
        ],
        'extended_vocab': [
            ('brown', '棕色'), ('gray', '灰色'), ('gold', '金色'),
            ('silver', '银色'), ('rectangle', '长方形'), ('star', '星形'),
            ('heart', '心形'), ('diamond', '菱形'),
        ],
        'video_search': 'English Singsing Color rainbow What color is it',
        'video_bv': 'BV1hvD2BkEDh',  # Collection with color video
        'video_file': 'colors_vocab.mp4',
        'video_file2': 'shapes_vocab.mp4',  # Shapes supplementary
    },
    {
        'num': 3,
        'title': 'Numbers & Time 数字与时间',
        'subtitle': "What time is it?",
        'icon': '🔢',
        'vocab': [
            ('one', '一'), ('two', '二'), ('three', '三'), ('four', '四'),
            ('five', '五'), ('six', '六'), ('seven', '七'), ('eight', '八'),
            ('nine', '九'), ('ten', '十'), ('clock', '时钟'), ('time', '时间'),
        ],
        'dialogue': [
            ('A', 'What time is it?', '几点了？'),
            ('B', "It's 8 o'clock. Time for breakfast!", '8点了。该吃早饭了！'),
            ('A', 'What time is it now?', '现在几点了？'),
            ('B', "It's 12 o'clock. Time for lunch!", '12点了。该吃午饭了！'),
            ('A', "I'm hungry! Let's eat!", '我饿了！我们吃饭吧！'),
            ('B', 'What time do you go to bed?', '你几点睡觉？'),
            ('A', 'I go to bed at 9 o\'clock.', '我9点睡觉。'),
        ],
        'grammar': [
            ('时间表达', "It's + 数字 + o'clock."),
            ('What time 问句', 'What time is it? / What time do you...?'),
        ],
        'extended_vocab': [
            ('eleven', '十一'), ('twelve', '十二'), ('twenty', '二十'),
            ('thirty', '三十'), ('hour', '小时'), ('minute', '分钟'),
            ('morning', '上午'), ('afternoon', '下午'),
        ],
        'video_search': 'English Singsing Number 123 What time is it',
        'video_bv': 'BV1BW4y1C7Cf',  # 480集 collection
        'video_file': 'numbers_vocab.mp4',
    },
    {
        'num': 4,
        'title': 'Body & Feelings 身体与感受',
        'subtitle': "How do you feel?",
        'icon': '😊',
        'vocab': [
            ('head', '头'), ('eyes', '眼睛'), ('nose', '鼻子'),
            ('mouth', '嘴巴'), ('ears', '耳朵'), ('hands', '手'),
            ('arms', '手臂'), ('legs', '腿'), ('feet', '脚'),
            ('happy', '开心'), ('sad', '难过'), ('angry', '生气'),
        ],
        'dialogue': [
            ('A', 'How do you feel today?', '你今天感觉怎么样？'),
            ('B', "I'm happy! The sun is shining!", '我很开心！阳光明媚！'),
            ('A', "That's great! I'm happy too.", '太棒了！我也很开心。'),
            ('B', 'Point to your head!', '指你的头！'),
            ('A', 'This is my head. These are my eyes.', '这是我的头。这是我的眼睛。'),
            ('B', 'Touch your nose! Clap your hands!', '摸你的鼻子！拍拍手！'),
            ('A', 'Stamp your feet! Well done!', '跺跺脚！做得好！'),
        ],
        'grammar': [
            ('感官动词', 'Touch / Point to / Clap / Stamp + 身体部位'),
            ('How 问句(感受)', 'How do you feel? → I\'m happy/sad/angry.'),
        ],
        'extended_vocab': [
            ('shoulders', '肩膀'), ('knees', '膝盖'), ('toes', '脚趾'),
            ('fingers', '手指'), ('tired', '累的'), ('excited', '兴奋的'),
            ('scared', '害怕的'), ('surprised', '惊讶的'),
        ],
        'video_search': 'English Singsing Body parts Feelings emotions',
        'video_bv': 'BV1f3411B7cx',  # 28主题词汇
        'video_file': 'body_vocab.mp4',
    },
    {
        'num': 5,
        'title': 'Family 家人',
        'subtitle': "Who is he? Who is she?",
        'icon': '👨‍👩‍👧‍👦',
        'vocab': [
            ('family', '家庭'), ('father', '爸爸'), ('mother', '妈妈'),
            ('brother', '兄弟'), ('sister', '姐妹'), ('baby', '宝宝'),
            ('grandpa', '爷爷'), ('grandma', '奶奶'), ('uncle', '叔叔'),
            ('aunt', '阿姨'), ('cousin', '表亲'), ('parents', '父母'),
        ],
        'dialogue': [
            ('A', 'Who is he?', '他是谁？'),
            ('B', 'He is my father.', '他是我爸爸。'),
            ('A', 'Who is she?', '她是谁？'),
            ('B', 'She is my mother.', '她是我妈妈。'),
            ('A', 'Do you have any brothers?', '你有兄弟吗？'),
            ('B', 'Yes, I have one brother.', '有，我有一个哥哥。'),
            ('A', 'I have a sister. We are a happy family!', '我有一个妹妹。我们是幸福的一家！'),
        ],
        'grammar': [
            ('Who 问句', 'Who is he/she? → He/She is my + 家庭成员.'),
            ('have 表达', 'I have + a/an + 家庭成员.'),
        ],
        'extended_vocab': [
            ('son', '儿子'), ('daughter', '女儿'), ('husband', '丈夫'),
            ('wife', '妻子'), ('nephew', '侄子'), ('niece', '侄女'),
            ('twins', '双胞胎'), ('relative', '亲戚'),
        ],
        'video_search': 'English Singsing Family members Who is he she',
        'video_bv': 'BV1aJ411d7Xm',  # 700+ collection, ep 13
        'video_file': 'family_vocab.mp4',
    },
    {
        'num': 6,
        'title': 'Food & Drinks 食物与饮品',
        'subtitle': "May I take your order?",
        'icon': '🍎',
        'vocab': [
            ('apple', '苹果'), ('banana', '香蕉'), ('bread', '面包'),
            ('milk', '牛奶'), ('water', '水'), ('juice', '果汁'),
            ('rice', '米饭'), ('chicken', '鸡肉'), ('pizza', '披萨'),
            ('cake', '蛋糕'), ('ice cream', '冰淇淋'), ('candy', '糖果'),
        ],
        'dialogue': [
            ('A', 'May I take your order?', '可以点餐了吗？'),
            ('B', "Yes, I'd like a pizza, please.", '好的，我想要一份披萨。'),
            ('A', 'Anything else?', '还要别的吗？'),
            ('B', 'A glass of juice, please.', '请来一杯果汁。'),
            ('A', 'For here or to go?', '在这吃还是带走？'),
            ('B', 'For here, please. Thank you!', '在这吃，谢谢！'),
            ('A', 'Do you like ice cream?', '你喜欢冰淇淋吗？'),
            ('B', 'Yes, I love it!', '是的，我超喜欢！'),
        ],
        'grammar': [
            ('点餐句型', "I'd like + 食物 / A glass of + 饮品."),
            ('选择疑问句', 'For here or to go?'),
            ('Do you like...?', 'Do you like + 食物? → Yes, I do / No, I don\'t.'),
        ],
        'extended_vocab': [
            ('hamburger', '汉堡'), ('sandwich', '三明治'), ('salad', '沙拉'),
            ('soup', '汤'), ('coffee', '咖啡'), ('tea', '茶'),
            ('noodles', '面条'), ('chocolate', '巧克力'),
        ],
        'video_search': 'English Singsing Food Fruits Vegetables Order',
        'video_bv': 'BV1BW4y1C7Cf',  # 480集 collection
        'video_file': 'food_vocab.mp4',
    },
    {
        'num': 7,
        'title': 'Animals 动物',
        'subtitle': "What's your favorite animal?",
        'icon': '🐶',
        'vocab': [
            ('dog', '狗'), ('cat', '猫'), ('bird', '鸟'),
            ('fish', '鱼'), ('rabbit', '兔子'), ('elephant', '大象'),
            ('lion', '狮子'), ('tiger', '老虎'), ('monkey', '猴子'),
            ('duck', '鸭子'), ('pig', '猪'), ('cow', '牛'),
        ],
        'dialogue': [
            ('A', "What's your favorite animal?", '你最喜欢什么动物？'),
            ('B', 'I like dogs. They are friendly!', '我喜欢狗。它们很友好！'),
            ('A', 'I like cats. They are cute!', '我喜欢猫。它们很可爱！'),
            ('B', 'Look at the elephant! It\'s so big!', '看那头大象！它好大！'),
            ('A', 'The monkey is jumping! So funny!', '猴子在跳！太有趣了！'),
            ('B', 'What sound does a cow make?', '牛怎么叫？'),
            ('A', 'Moo~ Moo~', '哞~ 哞~'),
        ],
        'grammar': [
            ('喜好表达', 'I like + 动物. They are + 形容词.'),
            ('描述特征', "It's + 形容词 (big/small/cute/fast)."),
        ],
        'extended_vocab': [
            ('horse', '马'), ('sheep', '羊'), ('frog', '青蛙'),
            ('bear', '熊'), ('panda', '熊猫'), ('penguin', '企鹅'),
            ('giraffe', '长颈鹿'), ('dolphin', '海豚'),
        ],
        'video_search': 'English Singsing Animals Sea Animals Animal Sounds',
        'video_bv': 'BV1f3411B7cx',  # 28主题词汇
        'video_file': 'animals_vocab.mp4',
    },
    {
        'num': 8,
        'title': 'Daily Routine 日常生活',
        'subtitle': "My Day - What do you do every day?",
        'icon': '⏰',
        'vocab': [
            ('wake up', '起床'), ('get up', '起来'), ('wash face', '洗脸'),
            ('brush teeth', '刷牙'), ('have breakfast', '吃早餐'),
            ('go to school', '上学'), ('study', '学习'), ('have lunch', '吃午饭'),
            ('go home', '回家'), ('do homework', '做作业'),
            ('have dinner', '吃晚饭'), ('go to bed', '睡觉'),
        ],
        'dialogue': [
            ('A', 'What time do you get up?', '你几点起床？'),
            ('B', 'I get up at 7 o\'clock.', '我7点起床。'),
            ('A', 'What do you do after school?', '放学后你做什么？'),
            ('B', 'I do my homework first.', '我先做作业。'),
            ('A', "Then I play with my friends. It's fun!", '然后和朋友玩。很有趣！'),
            ('B', 'What time do you go to bed?', '你几点睡觉？'),
            ('A', 'I go to bed at 9:30.', '我9:30睡觉。'),
        ],
        'grammar': [
            ('一般现在时', 'I get up / I go / I do — 日常动作'),
            ('时间表达', 'at + 时间 (at 7 o\'clock / at 9:30)'),
            ('顺序词', 'First... Then... After that...'),
        ],
        'extended_vocab': [
            ('take a shower', '洗澡'), ('get dressed', '穿衣服'),
            ('read books', '读书'), ('watch TV', '看电视'),
            ('play games', '玩游戏'), ('clean room', '打扫房间'),
            ('feed the pet', '喂宠物'), ('pack bag', '收拾书包'),
        ],
        'video_search': 'English Singsing My Day Daily Routine',
        'video_bv': 'BV1yB4WzZEwt',  # My Day - Daily Routine
        'video_file': 'daily_routine.mp4',
    },
    {
        'num': 9,
        'title': 'Weather & Seasons 天气与季节',
        'subtitle': "How's the weather?",
        'icon': '🌈',
        'vocab': [
            ('sunny', '晴天'), ('rainy', '下雨'), ('cloudy', '多云'),
            ('windy', '刮风'), ('snowy', '下雪'), ('hot', '热'),
            ('cold', '冷'), ('warm', '温暖'), ('cool', '凉爽'),
            ('spring', '春天'), ('summer', '夏天'),
            ('autumn', '秋天'), ('winter', '冬天'),
        ],
        'dialogue': [
            ('A', "How's the weather today?", '今天天气怎么样？'),
            ('B', "It's sunny and warm!", '晴天，很温暖！'),
            ('A', "Let's go to the park!", '我们去公园吧！'),
            ('B', "Great idea! I love sunny days.", '好主意！我喜欢晴天。'),
            ('A', "What's your favorite season?", '你最喜欢什么季节？'),
            ('B', 'I like summer. I can swim!', '我喜欢夏天。可以游泳！'),
            ('A', 'I like winter. I can make a snowman!', '我喜欢冬天。可以堆雪人！'),
        ],
        'grammar': [
            ('天气问句', "How's the weather? → It's + 天气词."),
            ('季节表达', 'I like + 季节. I can + 活动.'),
        ],
        'extended_vocab': [
            ('rainbow', '彩虹'), ('storm', '暴风雨'), ('foggy', '有雾'),
            ('freezing', '极冷'), ('lightning', '闪电'), ('thunder', '打雷'),
            ('umbrella', '雨伞'), ('temperature', '温度'),
        ],
        'video_search': "English Singsing Weather How's the weather",
        'video_bv': 'BV1V64y1F76f',  # Weather theme
        'video_file': 'weather_theme.mp4',
    },
    {
        'num': 10,
        'title': 'School & Hobbies 学校与爱好',
        'subtitle': "What do you like to do?",
        'icon': '📚',
        'vocab': [
            ('school', '学校'), ('teacher', '老师'), ('student', '学生'),
            ('classroom', '教室'), ('book', '书'), ('pen', '笔'),
            ('pencil', '铅笔'), ('bag', '书包'), ('desk', '课桌'),
            ('sing', '唱歌'), ('dance', '跳舞'), ('draw', '画画'),
        ],
        'dialogue': [
            ('A', 'What do you like to do?', '你喜欢做什么？'),
            ('B', 'I like to sing and dance!', '我喜欢唱歌跳舞！'),
            ('A', 'That sounds fun! I like to draw.', '听起来很有趣！我喜欢画画。'),
            ('B', 'Can you draw a cat?', '你能画一只猫吗？'),
            ('A', 'Yes, I can! Look!', '可以啊！看！'),
            ('B', 'Wow, it\'s beautiful! You are an artist!', '哇，好漂亮！你是艺术家！'),
            ('A', 'What subject do you like?', '你喜欢什么科目？'),
            ('B', 'I like English and Art!', '我喜欢英语和美术！'),
        ],
        'grammar': [
            ('情态动词 can', 'Can you + 动词? → Yes, I can / No, I can\'t.'),
            ('like to + 动词', 'I like to + 动词原形.'),
            ('What 问句', 'What subject do you like?'),
        ],
        'extended_vocab': [
            ('math', '数学'), ('science', '科学'), ('music', '音乐'),
            ('sports', '运动'), ('swimming', '游泳'), ('running', '跑步'),
            ('reading', '阅读'), ('writing', '写作'),
        ],
        'video_search': 'English Singsing School Hobbies Sports',
        'video_bv': 'BV1f3411B7cx',  # 28主题词汇 P15 - School subjects
        'video_file': 'school_subjects.mp4',
    },
]

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def hex_to_rgb(hex_str):
    """Convert hex color string to RGBColor"""
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def set_slide_bg(slide, color):
    """Set slide background to solid color"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, radius=None):
    """Add a rectangle shape, optionally with rounded corners"""
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # Adjust corner radius
        shape.adjustments[0] = radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if border_color:
        shape.line.color.rgb = border_color
        if border_width:
            shape.line.width = border_width
    else:
        shape.line.fill.background()

    return shape

def add_circle(slide, left, top, size, fill_color=None, border_color=None):
    """Add a circle/oval shape"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape

def add_star(slide, left, top, size, fill_color):
    """Add a decorative star"""
    shape = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_name=FONT_CN, font_size=18,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a textbox with formatted text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    # Set East Asian font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), font_name)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=16, default_color=BLACK):
    """Add a textbox with multiple formatted lines.
    Each line is (text, font_size, color, bold) - size/color/bold optional"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line, str):
            p.text = line
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
        else:
            p.text = line[0]
            p.font.size = Pt(line[1] if len(line) > 1 else default_size)
            p.font.color.rgb = line[2] if len(line) > 2 else default_color
            if len(line) > 3:
                p.font.bold = line[3]

        p.font.name = FONT_CN
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:eaTypeface'), FONT_CN)

    return txBox

def add_speech_bubble(slide, left, top, width, height, text, role=None, fill_color=None,
                      font_size=16, text_color=BLACK, pointer_left=True):
    """Add a speech bubble shape"""
    # Use rounded rectangle as simplified speech bubble
    shape = add_rect(slide, left, top, width, height,
                     fill_color=fill_color or hex_to_rgb('F5F5F5'),
                     border_color=hex_to_rgb('CCCCCC'), border_width=Pt(1), radius=0.05)

    # Add text
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    if role:
        run = p.add_run()
        run.text = f"{role}: "
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.name = FONT_EN

    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.name = FONT_CN

    return shape

def add_video_placeholder(slide, left, top, width, height, label="🎬 点击此处插入视频"):
    """Add a video placeholder shape when video file is not available"""
    shape = add_rect(slide, left, top, width, height,
                     fill_color=hex_to_rgb('1A1A2E'),
                     border_color=hex_to_rgb('444466'), border_width=Pt(2), radius=0.03)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_decorative_dots(slide, count=8, color=None):
    """Add scattered decorative circles"""
    import random
    for _ in range(count):
        x = Inches(random.uniform(0.5, 12.5))
        y = Inches(random.uniform(0.3, 7.0))
        size = Inches(random.uniform(0.08, 0.2))
        c = color or hex_to_rgb('CCCCCC')
        add_circle(slide, x, y, size, fill_color=c)

def add_bottom_bar(slide, color):
    """Add a decorative bottom bar"""
    add_rect(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5), fill_color=color)

def add_top_title_bar(slide, title_text, color):
    """Add a colorful top title bar"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), fill_color=color)
    add_textbox(slide, Inches(1), Inches(0.15), Inches(11), Inches(0.7),
                title_text, font_size=32, color=WHITE, bold=True,
                font_name=FONT_TITLE)

# ============================================================
# SLIDE BUILDERS
# ============================================================

def build_cover(prs):
    """Build the cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    primary = hex_to_rgb('156082')
    accent = hex_to_rgb('E97132')

    # Full background
    set_slide_bg(slide, hex_to_rgb('F0F8FF'))

    # Top colored band
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(2.5), fill_color=primary)

    # Decorative circles
    for i in range(6):
        color = hex_to_rgb(('E97132', '0F9ED5', '4EA72E', 'A02B93', 'FFB347', '7EC8E3')[i])
        x = Inches(2 + i * 1.8)
        add_circle(slide, x, Inches(1.0), Inches(0.6), fill_color=color)

    # Title
    add_textbox(slide, Inches(2), Inches(2.8), Inches(9), Inches(1.2),
                '🌟 English Singsing 🌟', font_size=48, color=primary, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    add_textbox(slide, Inches(2), Inches(3.9), Inches(9), Inches(0.8),
                '英语启蒙互动教学课件', font_size=36, color=accent, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Subtitle
    add_textbox(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.6),
                '看动画 · 学单词 · 练对话 · 角色扮演 · 词卡扩展',
                font_size=20, color=hex_to_rgb('666666'), alignment=PP_ALIGN.CENTER)

    # Feature boxes
    features = [
        ('🎬', '原版动画', 'English Singsing\n生动有趣'),
        ('📖', '场景对话', '真实语境\n自然习得'),
        ('🎭', '角色扮演', '师生互动\n模仿练习'),
        ('🃏', '词卡扩展', '主题词汇\n系统学习'),
    ]
    for i, (icon, title, desc) in enumerate(features):
        x = Inches(1.5 + i * 2.8)
        card = add_rect(slide, x, Inches(5.7), Inches(2.4), Inches(1.5),
                       fill_color=WHITE, border_color=hex_to_rgb('DDDDDD'),
                       border_width=Pt(1), radius=0.05)
        add_textbox(slide, x + Inches(0.1), Inches(5.8), Inches(2.2), Inches(0.4),
                   f'{icon} {title}', font_size=16, color=primary, bold=True,
                   alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(6.3), Inches(2.2), Inches(0.7),
                   desc, font_size=12, color=hex_to_rgb('888888'),
                   alignment=PP_ALIGN.CENTER)

    add_decorative_dots(slide, 20, hex_to_rgb('DDDDDD'))

def build_howto(prs):
    """Build the how-to-use slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb('156082')
    accent = hex_to_rgb('E97132')

    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, '📋 使用说明 — 五步教学法', primary)

    steps = [
        ('1', '看动画视频', '播放 English Singsing 原版动画，\n让学生沉浸在英语环境中', '🎬', accent),
        ('2', '学核心词汇', '从视频中提取关键词汇，\n用卡片形式展示和讲解', '📝', hex_to_rgb('0F9ED5')),
        ('3', '读场景对话', '摘录视频中的对话，\n逐句跟读模仿发音', '💬', hex_to_rgb('4EA72E')),
        ('4', '角色扮演', '分角色朗读对话，\n师生互动，模拟真实场景', '🎭', hex_to_rgb('A02B93')),
        ('5', '词卡扩展', '基于主题扩展更多词汇，\n制作词卡系统学习', '🃏', hex_to_rgb('E97132')),
    ]

    for i, (num, title, desc, icon, color) in enumerate(steps):
        y = Inches(1.3 + i * 1.15)
        # Number circle
        add_circle(slide, Inches(1.5), y + Inches(0.1), Inches(0.6), fill_color=color)
        add_textbox(slide, Inches(1.5), y + Inches(0.15), Inches(0.6), Inches(0.5),
                   num, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Content
        add_textbox(slide, Inches(2.5), y, Inches(1.5), Inches(0.5),
                   icon, font_size=28, color=color, bold=True)
        add_textbox(slide, Inches(3.5), y, Inches(2), Inches(0.4),
                   title, font_size=22, color=color, bold=True)
        add_textbox(slide, Inches(3.5), y + Inches(0.4), Inches(5), Inches(0.7),
                   desc, font_size=14, color=hex_to_rgb('666666'))

        # Arrow connector line
        if i < len(steps) - 1:
            add_rect(slide, Inches(1.75), y + Inches(0.7), Inches(0.05), Inches(0.65),
                    fill_color=hex_to_rgb('DDDDDD'))

    add_decorative_dots(slide, 10, hex_to_rgb('EEEEEE'))

def build_unit_cover(prs, unit):
    """Build a unit cover slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])
    accent = hex_to_rgb(UNIT_THEMES[unit['num']][2])
    secondary = hex_to_rgb(UNIT_THEMES[unit['num']][1])

    set_slide_bg(slide, WHITE)

    # Large decorative shape top-right
    add_circle(slide, Inches(9), Inches(-1), Inches(5), fill_color=primary)
    add_circle(slide, Inches(11), Inches(3), Inches(2.5), fill_color=secondary)

    # Unit number
    add_circle(slide, Inches(1), Inches(0.8), Inches(1.2), fill_color=accent)
    add_textbox(slide, Inches(1), Inches(0.95), Inches(1.2), Inches(0.9),
               f"Unit\n{unit['num']:02d}", font_size=28, color=WHITE, bold=True,
               alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, Inches(1), Inches(2.5), Inches(8), Inches(1.0),
               f"{unit['icon']} {unit['title']}", font_size=42, color=primary, bold=True,
               font_name=FONT_TITLE)

    add_textbox(slide, Inches(1), Inches(3.5), Inches(8), Inches(0.7),
               unit['subtitle'], font_size=24, color=hex_to_rgb('888888'),
               font_name=FONT_EN)

    # Step flow preview
    steps = ['🎬 看视频', '📝 学词汇', '💬 读对话', '🎭 角色扮演', '🃏 词卡扩展']
    for i, step in enumerate(steps):
        x = Inches(1 + i * 2.3)
        card = add_rect(slide, x, Inches(5.0), Inches(2.0), Inches(1.0),
                       fill_color=WHITE, border_color=hex_to_rgb('DDDDDD'),
                       border_width=Pt(1), radius=0.05)
        add_textbox(slide, x + Inches(0.1), Inches(5.15), Inches(1.8), Inches(0.7),
                   step, font_size=14, color=primary, bold=True, alignment=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            add_textbox(slide, x + Inches(2.0), Inches(5.3), Inches(0.3), Inches(0.4),
                       '→', font_size=14, color=accent, bold=True)

    add_decorative_dots(slide, 8, hex_to_rgb('EEEEEE'))

def build_video_slide(prs, unit):
    """Build a video slide with embedded video or placeholder"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])

    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, f"🎬 {unit['title']} — 看动画学英语", primary)

    # Video area (left side)
    video_left = Inches(0.5)
    video_top = Inches(1.3)
    video_w = Inches(8.0)
    video_h = Inches(5.0)

    # Try to find the video file - check unit's video_file first
    video_file = None
    vf_name = unit.get('video_file', '')
    if vf_name:
        candidate = VIDEO_DIR / vf_name
        if candidate.exists():
            video_file = str(candidate)

    # Fallback: search for any merged MP4
    if not video_file:
        for f in VIDEO_DIR.glob("*_*.mp4"):
            if '[00]' not in f.name and '[01]' not in f.name:
                video_file = str(f)
                break

    if video_file and os.path.exists(video_file):
        try:
            # Embed video
            movie_shape = slide.shapes.add_movie(
                video_file,
                video_left, video_top, video_w, video_h,
                poster_frame_image=None
            )
            # Add label
            add_textbox(slide, Inches(0.5), Inches(6.5), Inches(8), Inches(0.4),
                       f'📺 双击播放: {os.path.basename(video_file)}',
                       font_size=10, color=hex_to_rgb('999999'))
        except Exception as e:
            print(f"  Video embed failed for {vf_name}: {e}")
            add_video_placeholder(slide, video_left, video_top, video_w, video_h,
                                f"🎬 {unit['title']}\n视频文件: {vf_name}")
    else:
        add_video_placeholder(slide, video_left, video_top, video_w, video_h,
                            f"🎬 {unit['title']}\n视频文件: {vf_name or '待下载'}\n\n请运行 download_videos.py 下载视频")

    # Right side - pre-watching questions
    right_x = Inches(9.0)
    q_card = add_rect(slide, right_x, Inches(1.3), Inches(3.8), Inches(5.0),
                     fill_color=hex_to_rgb('F8F8F8'), border_color=hex_to_rgb('DDDDDD'),
                     border_width=Pt(1), radius=0.04)

    add_textbox(slide, right_x + Inches(0.3), Inches(1.5), Inches(3.2), Inches(0.4),
               '👀 看视频前想一想：', font_size=16, color=primary, bold=True)

    questions = [
        '1. 视频里有哪些人物？',
        '2. 他们在说什么？',
        '3. 你能听懂哪些单词？',
        '4. 模仿他们的语气说一说！',
    ]
    for i, q in enumerate(questions):
        add_textbox(slide, right_x + Inches(0.3), Inches(2.2 + i * 0.7), Inches(3.2), Inches(0.5),
                   q, font_size=13, color=hex_to_rgb('555555'))

    # Post-watching
    add_textbox(slide, right_x + Inches(0.3), Inches(4.8), Inches(3.2), Inches(0.4),
               '✅ 看视频后：', font_size=16, color=hex_to_rgb('4EA72E'), bold=True)

    post_questions = [
        '✓ 学到了哪些新单词？',
        '✓ 对话发生在什么场景？',
        '✓ 你能和老师对话吗？',
    ]
    for i, q in enumerate(post_questions):
        add_textbox(slide, right_x + Inches(0.3), Inches(5.2 + i * 0.5), Inches(3.2), Inches(0.4),
                   q, font_size=12, color=hex_to_rgb('777777'))

def build_vocab_slide(prs, unit):
    """Build a vocabulary slide with word cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])
    accent = hex_to_rgb(UNIT_THEMES[unit['num']][2])

    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, f"📝 {unit['title']} — 核心词汇 (Key Words)", primary)

    vocab = unit['vocab']
    rows, cols = 3, 4
    card_w = Inches(2.8)
    card_h = Inches(1.6)
    start_x = Inches(0.6)
    start_y = Inches(1.3)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)

    for idx, (en, cn) in enumerate(vocab):
        row = idx // cols
        col = idx % cols
        if row >= rows:
            break

        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Card background
        card = add_rect(slide, x, y, card_w, card_h,
                       fill_color=WHITE, border_color=hex_to_rgb('DDDDDD'),
                       border_width=Pt(1), radius=0.06)

        # Top accent stripe
        add_rect(slide, x + Inches(0.3), y + Inches(0.1), Inches(2.2), Inches(0.06),
                fill_color=accent)

        # English word
        add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.4), Inches(0.6),
                   en, font_size=22, color=primary, bold=True,
                   alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

        # Chinese meaning
        add_textbox(slide, x + Inches(0.2), y + Inches(0.95), Inches(2.4), Inches(0.5),
                   cn, font_size=14, color=hex_to_rgb('888888'),
                   alignment=PP_ALIGN.CENTER)

    # Note at bottom
    add_textbox(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
               '💡 老师提示：先读英文，再用中文解释。用手指卡片，让学生大声跟读！',
               font_size=12, color=hex_to_rgb('AAAAAA'))

def build_dialogue_slide(prs, unit):
    """Build a dialogue slide with speech bubbles"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])
    accent = hex_to_rgb(UNIT_THEMES[unit['num']][2])

    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, f"💬 {unit['title']} — 场景对话 (Dialogue)", primary)

    dialogue = unit['dialogue']
    # Display as a chat-like layout
    left_margin = Inches(1.5)
    right_margin = Inches(1.5)
    y = Inches(1.3)

    role_colors = {
        'A': (hex_to_rgb('E3F2FD'), hex_to_rgb('156082')),  # Blue bubble
        'B': (hex_to_rgb('FFF3E0'), hex_to_rgb('E65100')),  # Orange bubble
    }

    for i, (role, en, cn) in enumerate(dialogue):
        bg, text_color = role_colors.get(role, (hex_to_rgb('F5F5F5'), BLACK))
        is_left = (i % 2 == 0)  # Alternate left/right

        bubble_w = Inches(5.5)
        bubble_h = Inches(0.75)

        if is_left:
            x = left_margin
        else:
            x = Inches(13.333) - left_margin - bubble_w

        # Bubble
        bubble = add_rect(slide, x, y, bubble_w, bubble_h,
                        fill_color=bg, border_color=hex_to_rgb('DDDDDD'),
                        border_width=Pt(0.5), radius=0.08)

        # Role + English
        tf = bubble.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(5)

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{role}: "
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = text_color
        run.font.name = FONT_EN

        run2 = p.add_run()
        run2.text = en
        run2.font.size = Pt(14)
        run2.font.color.rgb = text_color
        run2.font.name = FONT_EN

        p2 = tf.add_paragraph()
        run3 = p2.add_run()
        run3.text = cn
        run3.font.size = Pt(11)
        run3.font.color.rgb = hex_to_rgb('999999')
        run3.font.name = FONT_CN

        y += Inches(0.85)

    add_decorative_dots(slide, 5, hex_to_rgb('EEEEEE'))

def build_readalong_slide(prs, unit):
    """Build a read-along practice slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])
    accent = hex_to_rgb(UNIT_THEMES[unit['num']][2])

    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, f"🔊 {unit['title']} — 跟读练习 (Listen & Repeat)", primary)

    dialogue = unit['dialogue']
    y = Inches(1.3)

    for i, (role, en, cn) in enumerate(dialogue):
        # Row background alternating
        bg = hex_to_rgb('FAFAFA') if i % 2 == 0 else WHITE
        add_rect(slide, Inches(1), y, Inches(11), Inches(0.65), fill_color=bg)

        # Role badge
        badge_color = primary if role == 'A' else accent
        badge = add_rect(slide, Inches(1.2), y + Inches(0.08), Inches(0.5), Inches(0.45),
                        fill_color=badge_color, radius=0.1)
        add_textbox(slide, Inches(1.2), y + Inches(0.1), Inches(0.5), Inches(0.4),
                   role, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # English
        add_textbox(slide, Inches(2.0), y + Inches(0.05), Inches(5), Inches(0.3),
                   en, font_size=16, color=BLACK, bold=True, font_name=FONT_EN)

        # Chinese
        add_textbox(slide, Inches(2.0), y + Inches(0.35), Inches(5), Inches(0.25),
                   cn, font_size=11, color=hex_to_rgb('AAAAAA'))

        # Pronunciation notes
        notes = get_pronunciation_notes(en)
        if notes:
            add_textbox(slide, Inches(7.5), y + Inches(0.1), Inches(4), Inches(0.4),
                       f'🔤 {notes}', font_size=11, color=hex_to_rgb('0F9ED5'))

        y += Inches(0.7)

def get_pronunciation_notes(text):
    """Simple pronunciation hints for common patterns"""
    hints = {
        'Good morning': '/gʊd ˈmɔːrnɪŋ/ — "古德 猫宁"',
        'How are you': '/haʊ ɑːr juː/ — "好啊油"',
        'Nice to meet you': '/naɪs tə miːt juː/ — 注意 meet 长音',
        'thank you': '/θæŋk juː/ — 咬舌音 th',
        "I'm fine": '/aɪm faɪn/ — 双元音 ai',
        "What's your name": '/wʌts jɔːr neɪm/ — what\'s 连读',
        'Hello': '/həˈloʊ/ — "哈楼"',
    }
    for key, hint in hints.items():
        if key.lower() in text.lower():
            return hint
    return None

def build_grammar_slide(prs, unit):
    """Build a grammar point slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])
    accent = hex_to_rgb(UNIT_THEMES[unit['num']][2])

    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, f"📖 {unit['title']} — 语法小课堂 (Grammar)", primary)

    grammar = unit['grammar']
    for i, (title, content) in enumerate(grammar):
        y = Inches(1.5 + i * 1.8)

        # Grammar card
        card = add_rect(slide, Inches(1.5), y, Inches(10), Inches(1.4),
                       fill_color=WHITE, border_color=primary, border_width=Pt(2), radius=0.04)

        # Title with accent background
        add_rect(slide, Inches(1.5), y, Inches(2.5), Inches(0.5), fill_color=primary)
        add_textbox(slide, Inches(1.7), y + Inches(0.05), Inches(2.1), Inches(0.4),
                   title, font_size=16, color=WHITE, bold=True)

        # Content
        add_textbox(slide, Inches(4.3), y + Inches(0.1), Inches(6.5), Inches(0.4),
                   content, font_size=20, color=BLACK, bold=True, font_name=FONT_EN)

        # Example from dialogue
        if i < len(unit['dialogue']):
            example_en = unit['dialogue'][i][1]
            add_textbox(slide, Inches(4.3), y + Inches(0.6), Inches(6.5), Inches(0.4),
                       f'例：{example_en}', font_size=14, color=accent, font_name=FONT_EN)

        # Bulb tip
        tips = [
            '💡 从视频对话中自然学习语法，不刻意死记硬背',
            '💡 先理解意思，再观察句型规律',
            '💡 用学到的句型自己造句！',
        ]
        if i < len(tips):
            add_textbox(slide, Inches(4.3), y + Inches(0.95), Inches(6.5), Inches(0.4),
                       tips[i], font_size=11, color=hex_to_rgb('AAAAAA'))

def build_roleplay_slide(prs, unit):
    """Build a role-play practice slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])
    accent = hex_to_rgb(UNIT_THEMES[unit['num']][2])

    set_slide_bg(slide, hex_to_rgb('FFFDF7'))  # Warm paper-like bg
    add_top_title_bar(slide, f"🎭 {unit['title']} — 角色扮演 (Role Play)", primary)

    # Left: Role A
    role_a_x = Inches(0.5)
    add_rect(slide, role_a_x, Inches(1.3), Inches(5.8), Inches(5.2),
            fill_color=hex_to_rgb('E3F2FD'), border_color=primary,
            border_width=Pt(2), radius=0.04)

    # Avatar placeholder
    add_circle(slide, Inches(2.5), Inches(1.5), Inches(1.2), fill_color=primary)
    add_textbox(slide, Inches(2.5), Inches(1.7), Inches(1.2), Inches(0.9),
               '👦\n角色A', font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(2.9), Inches(4.8), Inches(0.4),
               '🎤 你说：', font_size=16, color=primary, bold=True)

    # Right: Role B
    role_b_x = Inches(7.0)
    add_rect(slide, role_b_x, Inches(1.3), Inches(5.8), Inches(5.2),
            fill_color=hex_to_rgb('FFF3E0'), border_color=accent,
            border_width=Pt(2), radius=0.04)

    add_circle(slide, Inches(8.5), Inches(1.5), Inches(1.2), fill_color=accent)
    add_textbox(slide, Inches(8.5), Inches(1.7), Inches(1.2), Inches(0.9),
               '👧\n角色B', font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(7.5), Inches(2.9), Inches(4.8), Inches(0.4),
               '🎤 你说：', font_size=16, color=accent, bold=True)

    # Dialogue lines for each role
    dialogue = unit['dialogue']
    y_a = Inches(3.4)
    y_b = Inches(3.4)

    for i, (role, en, cn) in enumerate(dialogue):
        if role == 'A':
            add_textbox(slide, Inches(1), y_a, Inches(4.8), Inches(0.45),
                       f"{en}", font_size=15, color=BLACK, font_name=FONT_EN)
            add_textbox(slide, Inches(1), y_a + Inches(0.25), Inches(4.8), Inches(0.2),
                       f"{cn}", font_size=10, color=hex_to_rgb('999999'))
            y_a += Inches(0.5)
        else:
            add_textbox(slide, Inches(7.5), y_b, Inches(4.8), Inches(0.45),
                       f"{en}", font_size=15, color=BLACK, font_name=FONT_EN)
            add_textbox(slide, Inches(7.5), y_b + Inches(0.25), Inches(4.8), Inches(0.2),
                       f"{cn}", font_size=10, color=hex_to_rgb('999999'))
            y_b += Inches(0.5)

    # Bottom instruction
    add_textbox(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
               '🎯 老师和学生分角色朗读。先看着文字读，然后尝试不看文字凭记忆对话！可以交换角色再练一遍。',
               font_size=13, color=hex_to_rgb('888888'), alignment=PP_ALIGN.CENTER)

def build_wordcards_slide(prs, unit):
    """Build an extended word cards slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])
    accent = hex_to_rgb(UNIT_THEMES[unit['num']][2])

    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, f"🃏 {unit['title']} — 单词卡片墙 (Word Cards)", primary)

    words = unit['extended_vocab']
    cards_per_row = 4
    card_w = Inches(2.8)
    card_h = Inches(1.4)
    start_x = Inches(0.8)
    start_y = Inches(1.3)
    gap = Inches(0.3)

    for idx, (en, cn) in enumerate(words):
        row = idx // cards_per_row
        col = idx % cards_per_row

        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        # Flip card style
        card = add_rect(slide, x, y, card_w, card_h,
                       fill_color=WHITE, border_color=accent,
                       border_width=Pt(2), radius=0.05)

        # Top color bar
        add_rect(slide, x, y, card_w, Inches(0.08), fill_color=accent)

        # English
        add_textbox(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.5), Inches(0.5),
                   en, font_size=20, color=primary, bold=True,
                   alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

        # Chinese
        add_textbox(slide, x + Inches(0.15), y + Inches(0.75), Inches(2.5), Inches(0.4),
                   cn, font_size=13, color=hex_to_rgb('888888'),
                   alignment=PP_ALIGN.CENTER)

    # Bottom instruction
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
               '💡 老师提示：把卡片打印出来！正面写英文，背面写中文。每天抽5张复习，会的放一边，不会的继续练。',
               font_size=12, color=hex_to_rgb('AAAAAA'), alignment=PP_ALIGN.CENTER)

def build_exercise_slide(prs, unit):
    """Build a fun exercise/practice slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])
    accent = hex_to_rgb(UNIT_THEMES[unit['num']][2])

    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, f"🎮 {unit['title']} — 趣味练习 (Fun Practice)", primary)

    # Exercise 1: Match
    add_textbox(slide, Inches(1), Inches(1.3), Inches(5), Inches(0.4),
               '✏️ 练习一：连线 — 把英文和中文连起来', font_size=18, color=primary, bold=True)

    vocab = unit['vocab'][:6]
    # Left column: English
    for i, (en, cn) in enumerate(vocab):
        y = Inches(2.0 + i * 0.55)
        add_textbox(slide, Inches(1.5), y, Inches(3), Inches(0.4),
                   f'{i+1}. {en}', font_size=16, color=BLACK, font_name=FONT_EN)

    # Right column: Chinese (shuffled)
    import random
    cn_shuffled = random.sample([(i, cn) for i, (en, cn) in enumerate(vocab)], len(vocab))
    for j, (orig_idx, cn) in enumerate(cn_shuffled):
        y = Inches(2.0 + j * 0.55)
        add_textbox(slide, Inches(7), y, Inches(3), Inches(0.4),
                   f'{chr(65+j)}. {cn}', font_size=16, color=BLACK)

    add_textbox(slide, Inches(4), Inches(5.5), Inches(5), Inches(0.4),
               '答案：___  (学生连线后大声读出每个单词！)',
               font_size=12, color=hex_to_rgb('AAAAAA'), alignment=PP_ALIGN.CENTER)

    # Exercise 2: Fill in the blank
    add_textbox(slide, Inches(1), Inches(1.3), Inches(5), Inches(0.4),
               '', font_size=1)  # spacer moved

    # Second exercise on right side
    add_textbox(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.4),
               '✏️ 练习二：填空 — 补全对话', font_size=18, color=accent, bold=True)

    dialogue = unit['dialogue']
    if len(dialogue) >= 3:
        blank_line = dialogue[1][1]  # Take second line
        # Create blank by removing key words
        blank_text = blank_line.replace('morning', '______').replace('How', '____').replace('are', '___')
        add_textbox(slide, Inches(1.5), Inches(6.4), Inches(9), Inches(0.4),
                   f'A: Good ______! ____ ___ you?', font_size=18, color=BLACK,
                   alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

def build_summary_slide(prs, unit):
    """Build a unit summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])

    set_slide_bg(slide, hex_to_rgb('F0F8FF'))

    # Big checkmark
    add_textbox(slide, Inches(4), Inches(0.5), Inches(5), Inches(1.2),
               '✅', font_size=72, color=hex_to_rgb('4EA72E'), bold=True,
               alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(1.8), Inches(9), Inches(0.8),
               f'{unit["icon"]} Unit {unit["num"]:02d}: {unit["title"]} — 完成！',
               font_size=32, color=primary, bold=True, alignment=PP_ALIGN.CENTER,
               font_name=FONT_TITLE)

    # What we learned
    items = [
        f'✅ 掌握了 {len(unit["vocab"])} 个核心词汇',
        f'✅ 学会了 {len(unit["dialogue"])} 句日常对话',
        f'✅ 理解了 {len(unit["grammar"])} 个语法要点',
        f'✅ 完成了角色扮演练习',
        f'✅ 扩展了 {len(unit["extended_vocab"])} 个相关词汇',
    ]

    for i, item in enumerate(items):
        y = Inches(3.0 + i * 0.6)
        add_textbox(slide, Inches(3), y, Inches(7), Inches(0.5),
                   item, font_size=18, color=hex_to_rgb('444444'))

    # Self-check
    add_textbox(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.6),
               '🌟 自我检查：不看提示，你能说出所有单词和对话吗？',
               font_size=18, color=hex_to_rgb('E97132'), bold=True, alignment=PP_ALIGN.CENTER)

def build_final_slide(prs):
    """Build the final ending slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    primary = hex_to_rgb('156082')

    set_slide_bg(slide, primary)

    add_textbox(slide, Inches(2), Inches(1.5), Inches(9), Inches(1.2),
               '🌈 Keep Learning!', font_size=52, color=WHITE, bold=True,
               alignment=PP_ALIGN.CENTER, font_name=FONT_EN)

    add_textbox(slide, Inches(2), Inches(2.8), Inches(9), Inches(0.8),
               '坚持学习，每天进步一点点！', font_size=30, color=hex_to_rgb('7EC8E3'),
               alignment=PP_ALIGN.CENTER, font_name=FONT_TITLE)

    # Learning tips
    tips = [
        '📺 每天看一集 English Singsing 动画',
        '📝 用词卡复习学过的单词',
        '💬 大胆开口说英语，不怕犯错',
        '🎭 和家人朋友一起角色扮演',
        '🌟 坚持就是胜利！English is fun!',
    ]
    for i, tip in enumerate(tips):
        add_textbox(slide, Inches(3), Inches(3.8 + i * 0.6), Inches(7), Inches(0.5),
                   tip, font_size=16, color=WHITE)

    add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
               '📧 2075652822@qq.com  |  English Singsing © All videos from YouTube',
               font_size=10, color=hex_to_rgb('7EC8E3'), alignment=PP_ALIGN.CENTER)

# ============================================================
# MAIN
# ============================================================

def main():
    print("=" * 60)
    print("English Singsing 英语教学PPT生成器")
    print("=" * 60)

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("\n📝 Building slides...")

    # Cover
    print("  封面...")
    build_cover(prs)

    # How to use
    print("  使用说明...")
    build_howto(prs)

    # Unit list overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_title_bar(slide, '📚 课程目录 — 10个主题单元', hex_to_rgb('156082'))

    for i, unit in enumerate(UNITS):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.3 + row * 1.1)

        primary = hex_to_rgb(UNIT_THEMES[unit['num']][0])

        card = add_rect(slide, x, y, Inches(5.8), Inches(0.9),
                       fill_color=WHITE, border_color=primary, border_width=Pt(1.5), radius=0.04)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.35),
                   f'Unit {unit["num"]:02d}: {unit["icon"]} {unit["title"]}',
                   font_size=18, color=primary, bold=True)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.4), Inches(0.3),
                   unit['subtitle'], font_size=12, color=hex_to_rgb('888888'), font_name=FONT_EN)

    add_decorative_dots(slide, 8, hex_to_rgb('EEEEEE'))

    # Build each unit
    for unit in UNITS:
        print(f"  Unit {unit['num']:02d}: {unit['title']}...")

        build_unit_cover(prs, unit)        # 1. Unit cover
        build_video_slide(prs, unit)       # 2. Watch video
        build_vocab_slide(prs, unit)       # 3. Key vocabulary
        build_dialogue_slide(prs, unit)    # 4. Scene dialogue
        build_readalong_slide(prs, unit)   # 5. Read along
        build_grammar_slide(prs, unit)     # 6. Grammar
        build_roleplay_slide(prs, unit)    # 7. Role play
        build_wordcards_slide(prs, unit)   # 8. Word cards
        build_exercise_slide(prs, unit)    # 9. Fun practice
        # 10. Summary
        build_summary_slide(prs, unit)

    # Final slide
    print("  结束页...")
    build_final_slide(prs)

    # Save
    print(f"\n💾 Saving to: {OUTPUT_FILE}")
    prs.save(str(OUTPUT_FILE))

    slide_count = len(prs.slides)
    print(f"\n✅ Done! {slide_count} slides created.")
    print(f"📁 File: {OUTPUT_FILE}")
    print(f"📏 Size: {OUTPUT_FILE.stat().st_size / 1024 / 1024:.1f} MB")

if __name__ == '__main__':
    main()
